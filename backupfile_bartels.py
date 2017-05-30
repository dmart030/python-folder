from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from os import listdir, chdir
from os.path import isfile, join
import time, re, imghdr, glob, os, fnmatch
import numpy as np
import matplotlib.pyplot as plt
from PIL import Image


def image_edit(img3, i):
    if img3.find('PL.JPG') == -1 and img3.find('raman.JPG')== -1:
        if img3.find("edited") == -1:
            img = Image.open(img3, 'r')
            width, height = img.size
            img = img.crop(
                (
                    width - 1867,
                    height - 971,
                    width - 975,
                    height - 160
                    )
                )
            y = str(img3)
            y = y.strip("*.JPG")
            img.save(y + " edited image number " + str(i) + ".JPG")
            new_img_jpg[i] = (y + " edited image number " + str(i) + ".JPG")
    else:
        print('Error')
def command_1(k):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = Inches(2)
    pic = slide.shapes.add_picture(img_path_1, Inches(0), Inches(2), Inches(5), Inches(5) )
            
    left = Inches(5)
    top = Inches(2)
    height = Inches(5)
    pic = slide.shapes.add_picture(img_path_2, left, top, height=height, width = height)

def command_4(k, img):
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = img[:12] #CHANGE THIS IF THE NUMBERS INCREASE PAST 9999


    left = top = Inches(2)
    pic = slide.shapes.add_picture(img_path_3, Inches(0), Inches(2), Inches(5), Inches(5) )
    
def command_5(k):
    left = Inches(5)
    top = Inches(2)
    height = Inches(5)
    pic = slide.shapes.add_picture(img_path_4, left, top, height=height, width = height)

def command_6(k, y):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = Inches(2)
    pic = slide.shapes.add_picture(img_path_1, Inches(0), Inches(2), Inches(5), Inches(5) )
def numericalSort(value):
    parts = numbers.split(value)
    parts[1::2] = map(int, parts[1::2])
    return parts

def PL_measurement():
    i = 0
    while i!= x2:
        DataIn = np.loadtxt(all_txt_PL_list[i])
        y = list()
        v = list()
        for column in DataIn:
            v.append(column[1])
            y.append(column[0])
        plt.figure()
        plt.title('PL Spectrum')
        plt.ylabel('Intensity A.U')
        plt.xlabel('energy [eV]')
        DataIn_PL_name = str(all_txt_PL_list[i])
        graphs_PL.append('PL_plot_coordinates_' + DataIn_PL_name.strip('.txt')+ '.JPG')
        plt.plot(y,v, 'b')
        plt.savefig(DataIn_PL_name.strip('.txt') + '.JPG')
        y[:] = []
        v[:] = []
        i += 1
    
def RAMAN_measurement():
    j =0
    while j!= x0:
        DataIn_raman = np.loadtxt(all_txt_raman_list[j])
        y = list()
        v = list()
        for column in DataIn_raman:
            v.append(column[1])
            y.append(column[0])
        plt.figure()
        plt.ylabel('Intensity (A.U)')
        plt.title('Raman Spectrum')
        plt.xlabel('Raman shift $cm^{-1}$')
        plt.plot(y,v, 'b')
        DataIn_raman_name = str(all_txt_raman_list[j])
        graphs_raman.append('raman_plot_coordinates_' + DataIn_raman_name.strip('.txt') + 'raman.JPG')
        plt.savefig(DataIn_raman_name.strip('.txt') + '.JPG')
        y[:] = []
        v[:] = []
        j += 1
def image_find_jpg(x):
    #*********************************
    for file in glob.glob("*.JPG"):
        if file.find("PL.JPG") > -1 or file.find("raman.txt") > -1 :
            continue
        elif file.find("edited") == -1:
            x.append(file)
        else:
            continue
    x = sorted(x, key=numericalSort)
    #*********************************
    return x
def all_txt_PL(x):
    #*********************************
    for file in glob.glob("*PL.txt"):
        x.append(file)
    x = sorted(x, key=numericalSort)
    #*********************************
    return x
def all_txt_raman(x):
    #*********************************
    for file in glob.glob("*raman.txt"):
        x.append(file)
    x = sorted(x, key=numericalSort)
    #*********************************
    return x
def all_img_raman(x):
    #*********************************
    for file in glob.glob("*raman.jpg"):
        x.append(file)
    x = sorted(x, key=numericalSort)
    #*********************************
    return x
def all_img_PL(x):
    #*********************************
    for file in glob.glob("*PL.jpg"):
        x.append(file)
    x = sorted(x, key=numericalSort)
    #*********************************
    return x

numbers = re.compile(r'(\d+)')
inputthingy = input('what is the folder location of the files: ')
os.chdir(inputthingy) #what directory are the files
inputthingy1 = inputthingy + '"\"'


################title slide creation#################################################################################
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

 
title.text = "Dominic Martinez-Ta" #insert your name here
subtitle.text = "Created on: " + time.strftime("%c") ##you can choose to omit this if you don't want it.
######################################################################################################################

onlyfiles = list([i for i in listdir(inputthingy) if isfile(join(inputthingy, i))])

onlyfiles = sorted(onlyfiles, key=numericalSort)

x = len(onlyfiles) #measures the number of files in the folder

all_jpg = list()
all_txt = list()
all_txt_raman_list = list()
all_txt_PL_list = list()
cropped_images = list()

all_txt_PL_list = all_txt_PL(all_txt_PL_list)
all_txt_raman_list = all_txt_raman(all_txt_raman_list)

x0 = len(all_txt_raman_list)
x2 = len(all_txt_PL_list)
graphs_PL = list()
graphs_raman = list()


image_find_jpg(all_jpg)
x1 = len(all_jpg)
while True:
    print("Are these images already cropped?[yes/no]")
    response = input().lower()
    if response == 'yes':
        k = 0
        PL_measurement() #creates PL images
        RAMAN_measurement() #creates raman images
        image_find_jpg(new_img_jpg) #adds images to the empty list
        new_img_jpg = sorted(new_img_jpg, key=numericalSort) #orders the images in a proper way
        y = len(new_img_jpg)
        
        while k != y:
            if y % 2 == 0 and k < (y-1):
                img_path_1 = str(new_img_jpg[k])
                img_path_2 = str(new_img_jpg[k+1])
                command_1(k) #creates slide of two images
                if (k < (y-2)):
                    img_path_3 = str(new_img_jpg[k+2])
                    if img_path_3.find('PL') > -1:
                        command_4(k) #creates slide of one image
                        if x0 > 0:
                            img_path_4 = str(new_img_jpg[k+3])
                            if img_path_4.find("raman.JPG") > -1:
                                  command_5(k) #adds to slide of one image
                                  k += 1
                        k += 1
                k += 2
            elif k == (y-1):
                img_path_1 = str(new_img_jpg[k])
                command_6(k, new_img_jpg)
                k  += 1
            else:
                img_path_1 = str(new_img_jpg[k])
                img_path_2 = str(new_img_jpg[k+1])
                command_1(k, new_img_jpg) #creates slide of two images
                if (x2 > 0) and (k < (y-2)):
                    img_path_3 = str(new_img_jpg[k+2])
                    if img_path_3.find('PL') > -1:
                        command_4(k,img_path_3)#creates slide of one image
                        if x0 > 0:
                            img_path_4 = str(new_img_jpg[k+3])
                            if img_path_4.find("raman.JPG") > -1:
                                command_5(k) #adds to slide of one image
                                k += 1
                        k += 1
                k += 2
        print(r"That's all for now, your powerpoint has been created :)")
        prs.save('group_update_.pptx')
        break
    elif response == 'no':
        all_jpg = sorted(all_jpg, key=numericalSort) #sorts the list in a logical manner
        i = 0
        l = 0
        j = 0
        k = 0
        new_img_jpg = [l for l in range(len(all_jpg))]
        for i in range(x1):
            image_edit(all_jpg[i], i) #edites each image individuall foudn in the all_jpg list
        PL_measurement() #creates PL images
        RAMAN_measurement() #creates raman images
        PL_raman_jpg_img_list = list() #creates an empty list obviousl
        all_img_raman(PL_raman_jpg_img_list) #adds raman images to empty list
        all_img_PL(PL_raman_jpg_img_list) #adds PL images to empty list
        i = 0
        for i in range(len(PL_raman_jpg_img_list)):
            all_jpg.append(PL_raman_jpg_img_list[i-1]) #appends raman and PL images to all_jpg list
        all_jpg = sorted(all_jpg, key=numericalSort)   
        for j in range(len(PL_raman_jpg_img_list)):
            new_img_jpg.append(PL_raman_jpg_img_list[j-1]) #appends raman and PL images to new_img_jpg list
        new_img_jpg = sorted(new_img_jpg, key=numericalSort)
        x2 = len(PL_raman_jpg_img_list) #determines the number of images for raman and PL
        y = len(new_img_jpg) #determines the length of new_img_jpg
        while k != y:
            if y % 2 == 0 and k < (y-1):
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                title_of_file_1 = all_jpg[k] #this created a title slide per two images. It's main job is to just say what samples these are.
                if (k+1) < y:
                    title_of_file_2 = all_jpg[k+1]
                    title_of_file_2 = title_of_file_2.strip('.JPG')
                    title_of_file_1 = title_of_file_1.strip('.JPG')
                if (k+1) < y:
                    title.text = title_of_file_1 + " and " + title_of_file_2
                else:
                    title.text = title_of_file_1
                img_path_1 = str(new_img_jpg[k])
                img_path_2 = str(new_img_jpg[k+1])
                command_1(k) #this purpose is to put the images on the slide
                if (k < (y-2)):
                    img_path_3 = str(new_img_jpg[k+2])
                    if img_path_3.find('PL') > -1: 
                        command_4(k, img_path_3)
                        if x0 > 0:
                            img_path_4 = str(new_img_jpg[k+3])
                            if img_path_4.find("raman.JPG") > -1:
                                command_5(k)
                                k += 1
                        k += 1
                k += 2
            elif k == (y-1):
                img_path_1 = str(new_img_jpg[k])
                command_6(k, new_img_jpg)
                k  += 1
            else:
                img_path_1 = str(new_img_jpg[k])
                img_path_2 = str(new_img_jpg[k+1])
                command_1(k)
                if (k < (y-2)):
                    img_path_3 = str(new_img_jpg[k+2])
                    if img_path_3.find('PL.JPG') > -1: 
                        command_4(k, img_path_3)
                        if x0 > 0:
                            img_path_4 = str(new_img_jpg[k+1])
                            if img_path_4.find("raman.JPG") > -1:
                                command_5(k)
                                k += 1
                        k += 1
                k += 2
        print(r"That's all for now, your powerpoint has been created :)")
        prs.save('group_update_.pptx')
        break
    else:
        print('ERROR \n Please input a yes or no response.')
        continue



