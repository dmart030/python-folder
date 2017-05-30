from pptx import Presentation
from pptx.util import Inches

img_path = 'WIN_20160628_09_54_25_Pro.jpg'

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(2)
pic = slide.shapes.add_picture(img_path, Inches(0), Inches(2), Inches(5), Inches(5) )

left = Inches(5)
top = Inches(2)
height = Inches(5)
pic = slide.shapes.add_picture(img_path, left, top, height=height, width = height)

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)


pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(7), Inches(7) )

prs.save('inputing_an_image.pptx')
